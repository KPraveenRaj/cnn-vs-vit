"""Shared PowerPoint helpers. House style lives here and nowhere else.

Matches the style already established by the submitted proposal deck
(report/phase1/proposal/make_ppt.py): 16:9, dark-blue accent #1F4E79, near-black
body text. Model colours are the two categorical slots used by every figure, so
a colour means the same thing on a slide as it does inside a plot.

Layout constants are named rather than inlined so a slide reads as content
instead of geometry, and so the whole set restyles from one edit.
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ACCENT = RGBColor(0x1F, 0x4E, 0x79)
BODY = RGBColor(0x26, 0x26, 0x26)
MUTED = RGBColor(0x59, 0x59, 0x59)
RULE = RGBColor(0xD9, 0xD9, 0xD6)
GOOD = RGBColor(0x1B, 0xAF, 0x7A)
RESNET = RGBColor(0x2A, 0x78, 0xD6)
VIT = RGBColor(0xEB, 0x68, 0x34)

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.62)                      # page margin
CONTENT_W = W - 2 * M


def new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _tb(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    return box.text_frame


def title_slide(prs, title, subtitle, footer=None):
    s = _blank(prs)
    tf = _tb(s, M, Inches(2.15), CONTENT_W, Inches(2.0))
    p = tf.paragraphs[0]
    p.text = title
    p.font.size, p.font.bold, p.font.color.rgb = Pt(40), True, ACCENT
    p = tf.add_paragraph()
    p.text = subtitle
    p.font.size, p.font.color.rgb = Pt(19), BODY
    p.space_before = Pt(14)
    if footer:
        tf2 = _tb(s, M, Inches(6.05), CONTENT_W, Inches(1.0))
        for i, line in enumerate(footer):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p.text = line
            p.font.size, p.font.color.rgb = Pt(13), MUTED
    return s


def slide(prs, title, kicker=None):
    """Standard content slide: accent title, optional small kicker above it."""
    s = _blank(prs)
    top = Inches(0.34)
    if kicker:
        tf = _tb(s, M, top, CONTENT_W, Inches(0.32))
        p = tf.paragraphs[0]
        p.text = kicker.upper()
        p.font.size, p.font.bold, p.font.color.rgb = Pt(11), True, MUTED
        top = Inches(0.66)
    tf = _tb(s, M, top, CONTENT_W, Inches(0.75))
    p = tf.paragraphs[0]
    p.text = title
    p.font.size, p.font.bold, p.font.color.rgb = Pt(27), True, ACCENT
    ln = s.shapes.add_shape(1, M, Inches(1.34), CONTENT_W, Emu(9525))
    ln.fill.solid(); ln.fill.fore_color.rgb = RULE
    ln.line.fill.background(); ln.shadow.inherit = False
    return s


def bullets(slide_, items, top=Inches(1.62), left=M, width=None, size=17,
            height=Inches(5.2)):
    """items: (level, text, bold) or plain str (level 0, not bold)."""
    tf = _tb(slide_, left, top, width or CONTENT_W, height)
    first = True
    for it in items:
        level, text, bold = (0, it, False) if isinstance(it, str) else it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = ("• " if level == 0 else "– ") + text if text else ""
        p.level = level
        p.font.size = Pt(size - 2 * level)
        p.font.bold = bold
        p.font.color.rgb = BODY if level == 0 else MUTED
        p.space_after = Pt(9)
    return tf


def picture(slide_, path, top=Inches(1.6), max_h=Inches(5.2), left=None, max_w=None):
    """Insert an image scaled to fit, horizontally centred by default."""
    from PIL import Image
    max_w = max_w or CONTENT_W
    with Image.open(path) as im:
        ar = im.width / im.height
    h, w = max_h, int(max_h * ar)
    if w > max_w:
        w, h = max_w, int(max_w / ar)
    left = left if left is not None else int(M + (CONTENT_W - w) / 2)
    return slide_.shapes.add_picture(path, left, top, width=w, height=h)


def table(slide_, headers, rows, top=Inches(1.7), left=M, width=None,
          height=None, size=13, col_w=None, highlight_rows=()):
    width = width or CONTENT_W
    height = height or Inches(0.42 * (len(rows) + 1))
    shape = slide_.shapes.add_table(len(rows) + 1, len(headers), left, top,
                                    width, height)
    tbl = shape.table
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Emu(int(width * cw / total))
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = str(h)
        p = cell.text_frame.paragraphs[0]
        p.font.size, p.font.bold, p.font.color.rgb = Pt(size), True, RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
    for r, row in enumerate(rows, start=1):
        for c, v in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(v)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(size)
            p.font.color.rgb = BODY
            p.font.bold = r - 1 in highlight_rows
            cell.fill.solid()
            cell.fill.fore_color.rgb = (RGBColor(0xEC, 0xF3, 0xFA) if r - 1 in highlight_rows
                                        else RGBColor(0xFF, 0xFF, 0xFF))
    return tbl


def caption(slide_, text, top=Inches(6.75), size=11):
    tf = _tb(slide_, M, top, CONTENT_W, Inches(0.6))
    p = tf.paragraphs[0]
    p.text = text
    p.font.size, p.font.color.rgb, p.font.italic = Pt(size), MUTED, True
    return tf


def takeaway(slide_, text, top=Inches(6.35)):
    """One boxed sentence: what the examiner should remember from this slide."""
    box = slide_.shapes.add_shape(5, M, top, CONTENT_W, Inches(0.72))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xEC, 0xF3, 0xFA)
    box.line.color.rgb = RULE
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size, p.font.bold, p.font.color.rgb = Pt(14), True, ACCENT
    p.alignment = PP_ALIGN.LEFT
    return box


def save(prs, path):
    from pathlib import Path
    Path(path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    n = len(prs.slides.__iter__.__self__._sldIdLst)
    print(f"  {Path(path).name}  ({n} slides)")


def notes(slide_, text):
    """Attach speaker notes — the script for what to SAY on this slide.

    These carry most of the accessibility burden. The slide itself stays sparse
    enough to project; the notes explain every term the first time it appears and
    give a plain-language reading of each number, so the deck can be presented to
    an audience with no machine-learning background without the presenter having
    to improvise the explanation.
    """
    tf = slide_.notes_slide.notes_text_frame
    first = True
    for para_text in [t.strip() for t in text.strip().split("\n\n") if t.strip()]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = para_text
        p.font.size = Pt(12)
    return slide_


def plain(slide_, text, top=Inches(5.55), height=Inches(1.05)):
    """A 'in plain terms' box: the same claim, stripped of jargon.

    Sits on the slide itself rather than in the notes, because a non-expert in
    the room needs it visible while they look at the figure, not afterwards.
    """
    box = slide_.shapes.add_shape(5, M, top, CONTENT_W, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xF6, 0xE8)
    box.line.color.rgb = RGBColor(0xE8, 0xC9, 0x9A)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = "In plain terms"
    p.font.size, p.font.bold, p.font.color.rgb = Pt(11), True, RGBColor(0x9A, 0x6B, 0x1E)
    p2 = tf.add_paragraph()
    p2.text = text
    p2.font.size, p2.font.color.rgb = Pt(13.5), BODY
    return box


def two_col(slide_, left_items, right_items, top=Inches(1.62), size=15,
            height=Inches(4.6)):
    """Two bullet columns — for side-by-side comparisons like CNN vs ViT."""
    w = int((CONTENT_W - Inches(0.4)) / 2)
    a = bullets(slide_, left_items, top=top, left=M, width=w, size=size, height=height)
    b = bullets(slide_, right_items, top=top, left=M + w + Inches(0.4), width=w,
                size=size, height=height)
    return a, b


def col_header(slide_, text, left_half=True, top=Inches(1.28), color=None):
    w = int((CONTENT_W - Inches(0.4)) / 2)
    left = M if left_half else M + w + Inches(0.4)
    tf = _tb(slide_, left, top, w, Inches(0.34))
    p = tf.paragraphs[0]
    p.text = text
    p.font.size, p.font.bold = Pt(15), True
    p.font.color.rgb = color or ACCENT
    return tf
